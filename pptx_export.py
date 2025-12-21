"""
Модуль експорту звіту у формат PowerPoint (.pptx).
ВЕРСІЯ: FINAL LOGO FIX (PNG).
- Використовує logo.png (SVG часто не підтримується).
- Абсолютний шлях до файлу (щоб точно знайти картинку).
- Кастомний останній слайд.
"""

import io
import os
import textwrap
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor

from classification import QuestionType
from summary import QuestionSummary

# --- НАЛАШТУВАННЯ ---
CHART_DPI = 150
FONT_SIZE_CHART = 10
BAR_WIDTH = 0.6
UNIV_NAME = "Національний університет імені Тараса Шевченка"

# ВАЖЛИВО: Використовуйте PNG, бо python-pptx не любить SVG
LOGO_FILENAME = "logo2.png" 
LOGO_PATH = os.path.join(os.getcwd(), LOGO_FILENAME)

def create_chart_image(qs: QuestionSummary) -> io.BytesIO:
    plt.close('all')
    plt.clf()
    plt.rcParams.update({'font.size': FONT_SIZE_CHART})
    
    labels = qs.table["Варіант відповіді"].astype(str).tolist()
    values = qs.table["Кількість"]
    wrapped_labels = [textwrap.fill(l, 25) for l in labels]

    is_scale = (qs.question.qtype == QuestionType.SCALE)
    if not is_scale:
        try:
            vals = pd.to_numeric(qs.table["Варіант відповіді"], errors='coerce')
            if vals.notna().all() and vals.min() >= 0 and vals.max() <= 10:
                is_scale = True
        except: pass

    if is_scale:
        fig = plt.figure(figsize=(6.0, 4.0))
        bars = plt.bar(wrapped_labels, values, color='#4F81BD', width=BAR_WIDTH)
        plt.ylabel('Кількість')
        plt.grid(axis='y', linestyle='--', alpha=0.5)
        for bar in bars:
            height = bar.get_height()
            plt.text(bar.get_x() + bar.get_width()/2., height + 0.1,
                     f'{int(height)}', ha='center', va='bottom', fontweight='bold')
    else:
        fig = plt.figure(figsize=(6.0, 4.0))
        colors = ['#4F81BD', '#C0504D', '#9BBB59', '#8064A2', '#4BACC6', '#F79646']
        c_arg = colors[:len(values)] if len(values) <= len(colors) else None
        
        wedges, texts, autotexts = plt.pie(
            values, labels=None, autopct='%1.1f%%', startangle=90,
            pctdistance=0.8, colors=c_arg, radius=1.0,
            textprops={'fontsize': FONT_SIZE_CHART}
        )
        for autotext in autotexts:
            autotext.set_color('white')
            autotext.set_weight('bold')
            import matplotlib.patheffects as path_effects
            autotext.set_path_effects([path_effects.withStroke(linewidth=2, foreground='#333333')])

        plt.axis('equal')
        cols = 2 if len(labels) > 3 else 1
        plt.legend(wrapped_labels, loc="upper center", bbox_to_anchor=(0.5, 0.0), ncol=cols, frameon=False, fontsize=9)

    plt.tight_layout()
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', dpi=CHART_DPI, bbox_inches='tight')
    plt.close(fig)
    img_stream.seek(0)
    return img_stream

def build_pptx_report(original_df, sliced_df, summaries, range_info) -> bytes:
    prs = Presentation()
    
    # === 1. ТИТУЛЬНИЙ СЛАЙД ===
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # ЛОГОТИП (Перевірка наявності файлу)
    if os.path.exists(LOGO_PATH):
        try:
            slide.shapes.add_picture(LOGO_PATH, Inches(0.2), Inches(0.2), width=Inches(1.2))
        except Exception as e:
            print(f"Помилка додавання лого: {e}")

    # Університет
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(0.4), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = UNIV_NAME
    p.font.bold = True
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.LEFT

    # Заголовки
    title = slide.shapes.title
    title.text = "Звіт про результати опитування"
    
    subtitle = slide.placeholders[1]
    subtitle.text = f"Всього анкет: {len(original_df)}\nОброблено: {len(sliced_df)}\nДіапазон: {range_info}"

    # === 2. СЛАЙДИ З ПИТАННЯМИ ===
    for qs in summaries:
        if qs.table.empty: continue
        
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = f"{qs.question.code}. {qs.question.text}"
        title_shape.text_frame.paragraphs[0].font.size = Pt(24)
        
        body = slide.placeholders[1]
        body.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        text_content = ""
        for row in qs.table.itertuples(index=False):
            val = str(row[0])
            count = str(row[1])
            pct = str(row[2])
            text_content += f"{val}: {count} ({pct}%)\n"
            
        body.text = text_content
        for paragraph in body.text_frame.paragraphs:
            paragraph.font.size = Pt(14)

        try:
            img_stream = create_chart_image(qs)
            slide.shapes.add_picture(img_stream, Inches(5), Inches(2), width=Inches(4.5))
        except: pass

    # === 3. ФІНАЛЬНИЙ СЛАЙД ===
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)

    # ЛОГОТИП (По центру зверху)
    if os.path.exists(LOGO_PATH):
        try:
            # Центруємо: (10 - 1.5) / 2 = 4.25
            slide.shapes.add_picture(LOGO_PATH, Inches(4.25), Inches(1.0), width=Inches(1.5))
        except: pass

    # ДЯКУЮ
    tb_thanks = slide.shapes.add_textbox(Inches(0), Inches(3.0), Inches(10), Inches(1.5))
    tf_thanks = tb_thanks.text_frame
    p_thanks = tf_thanks.paragraphs[0]
    p_thanks.text = "Дякую за увагу!"
    p_thanks.alignment = PP_ALIGN.CENTER
    p_thanks.font.size = Pt(32)
    p_thanks.font.bold = True
    p_thanks.font.color.rgb = RGBColor(0, 0, 0)

    # ІНФОРМАЦІЯ
    tb_info = slide.shapes.add_textbox(Inches(0), Inches(4.5), Inches(10), Inches(2))
    tf_info = tb_info.text_frame
    
    lines = [
        "Створено за допомогою додатку студентки МПУіК – Каптар Діани.",
        "Керівник проєкту – доцент Фратавчан Валерій Григорович."
    ]

    for line in lines:
        p = tf_info.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(80, 80, 80)
        p.space_after = Pt(10)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()